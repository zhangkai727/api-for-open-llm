""" from https://github.com/chatchat-space/Langchain-Chatchat """
import csv
import os
from io import BytesIO, TextIOWrapper
from pathlib import Path
from typing import (
    Dict,
    List,
    Optional,
    Sequence,
    TYPE_CHECKING,
)

import nltk
import numpy as np
from PIL import Image
from langchain.docstore.document import Document
from langchain.document_loaders.helpers import detect_file_encodings
from langchain_community.document_loaders.csv_loader import CSVLoader
from langchain_community.document_loaders.unstructured import UnstructuredFileLoader
from tqdm import tqdm

<<<<<<< HEAD
if TYPE_CHECKING:  # 如果进行类型检查
    try:
        from rapidocr_paddle import RapidOCR  # 尝试导入 rapidocr_paddle 包中的 RapidOCR 类
    except ImportError:
        from rapidocr_onnxruntime import RapidOCR  # 如果导入失败，则从 rapidocr_onnxruntime 包中导入 RapidOCR 类

PDF_OCR_THRESHOLD = (0.6, 0.6)  # 定义 PDF 光学字符识别的阈值元组 (文本分辨率阈值, 图像分辨率阈值)

NLTK_DATA_PATH = os.path.join(Path(__file__).parents[3], "assets", "nltk_data")  # 构建 NLTK 数据路径，假设该文件路径为从当前文件往上三级目录下的 assets/nltk_data 目录
nltk.data.path = [NLTK_DATA_PATH] + nltk.data.path  # 将 NLTK 数据路径添加到 NLTK 的搜索路径列表的开头

=======
if TYPE_CHECKING:
    try:
        from rapidocr_paddle import RapidOCR
    except ImportError:
        from rapidocr_onnxruntime import RapidOCR


PDF_OCR_THRESHOLD = (0.6, 0.6)

NLTK_DATA_PATH = os.path.join(Path(__file__).parents[3], "assets", "nltk_data")
nltk.data.path = [NLTK_DATA_PATH] + nltk.data.path
>>>>>>> d45db7c71cc1d7c6f454aab8dc32da6b0299ee3d


__all__ = [
    "FilteredCSVLoader",
    "RapidOCRDocLoader",
    "RapidOCRImageLoader",
    "RapidOCRPDFLoader",
    "OpenParserPDFLoader",
    "RapidOCRPPTLoader",
]


def get_ocr(use_cuda: bool = True) -> "RapidOCR":
    try:
<<<<<<< HEAD
        from rapidocr_paddle import RapidOCR  # 尝试从 rapidocr_paddle 包导入 RapidOCR 类
        # 使用 RapidOCR 类初始化 OCR 实例，指定是否在 GPU 上使用 CUDA 加速
        ocr = RapidOCR(det_use_cuda=use_cuda, cls_use_cuda=use_cuda, rec_use_cuda=use_cuda)
    except ImportError:
        from rapidocr_onnxruntime import RapidOCR  # 如果导入失败，则从 rapidocr_onnxruntime 包导入 RapidOCR 类
        ocr = RapidOCR()  # 使用默认配置初始化 OCR 实例
    return ocr  # 返回初始化后的 OCR 实例

=======
        from rapidocr_paddle import RapidOCR
        ocr = RapidOCR(det_use_cuda=use_cuda, cls_use_cuda=use_cuda, rec_use_cuda=use_cuda)
    except ImportError:
        from rapidocr_onnxruntime import RapidOCR
        ocr = RapidOCR()
    return ocr
>>>>>>> d45db7c71cc1d7c6f454aab8dc32da6b0299ee3d


class FilteredCSVLoader(CSVLoader):
    def __init__(
        self,
        file_path: str,
        columns_to_read: List[str],
        source_column: Optional[str] = None,
        metadata_columns: Sequence[str] = (),
        csv_args: Optional[Dict] = None,
        encoding: Optional[str] = None,
        autodetect_encoding: bool = False,
    ):
        super().__init__(
<<<<<<< HEAD
            file_path=file_path,  # 调用父类 CSVLoader 的构造函数，传入文件路径
            source_column=source_column,  # 指定源列名
            metadata_columns=metadata_columns,  # 指定元数据列名列表
            csv_args=csv_args,  # CSV 文件读取的参数字典
            encoding=encoding,  # 文件编码
            autodetect_encoding=autodetect_encoding,  # 是否自动检测文件编码
        )
        self.columns_to_read = columns_to_read  # 初始化时指定要读取的列名列表
=======
            file_path=file_path,
            source_column=source_column,
            metadata_columns=metadata_columns,
            csv_args=csv_args,
            encoding=encoding,
            autodetect_encoding=autodetect_encoding,
        )
        self.columns_to_read = columns_to_read
>>>>>>> d45db7c71cc1d7c6f454aab8dc32da6b0299ee3d

    def load(self) -> List[Document]:
        """Load data into document objects."""
        docs = []
        try:
<<<<<<< HEAD
            # 尝试打开 CSV 文件进行读取，使用指定的编码（如果有）
            with open(self.file_path, newline="", encoding=self.encoding) as csvfile:
                docs = self.__read_file(csvfile)  # 调用私有方法 __read_file 加载数据
        except UnicodeDecodeError as e:
            # 处理 Unicode 解码错误
            if self.autodetect_encoding:
                # 如果自动检测编码，则尝试多种编码进行加载
=======
            with open(self.file_path, newline="", encoding=self.encoding) as csvfile:
                docs = self.__read_file(csvfile)
        except UnicodeDecodeError as e:
            if self.autodetect_encoding:
>>>>>>> d45db7c71cc1d7c6f454aab8dc32da6b0299ee3d
                detected_encodings = detect_file_encodings(self.file_path)
                for encoding in detected_encodings:
                    try:
                        with open(
<<<<<<< HEAD
                                self.file_path, newline="", encoding=encoding.encoding
                        ) as csvfile:
                            docs = self.__read_file(csvfile)  # 使用检测到的编码加载数据
=======
                            self.file_path, newline="", encoding=encoding.encoding
                        ) as csvfile:
                            docs = self.__read_file(csvfile)
>>>>>>> d45db7c71cc1d7c6f454aab8dc32da6b0299ee3d
                            break
                    except UnicodeDecodeError:
                        continue
            else:
<<<<<<< HEAD
                # 如果不自动检测编码，则抛出 RuntimeError 异常
                raise RuntimeError(f"Error loading {self.file_path}") from e  # 抛出加载错误异常
        except Exception as e:
            # 捕获其他异常情况，如文件打开失败等，抛出 RuntimeError 异常
            raise RuntimeError(f"Error loading {self.file_path}") from e  # 抛出加载错误异常

        return docs  # 返回加载后的文档对象列表

    def __read_file(self, csvfile: TextIOWrapper) -> List[Document]:
        docs = []  # 初始化空列表，用于存储读取后的文档对象
        csv_reader = csv.DictReader(csvfile,
                                    **self.csv_args)  # 使用 csv.DictReader 读取 CSV 文件内容，**self.csv_args 将额外的参数传递给 DictReader

        # 遍历 CSV 文件的每一行
        for i, row in enumerate(csv_reader):
            content = []  # 初始化空列表，用于存储每行的内容

            # 遍历需要读取的列
            for col in self.columns_to_read:
                if col in row:
                    content.append(f'{col}:{str(row[col])}')  # 将列名和对应的值拼接成字符串，添加到 content 列表中
                else:
                    raise ValueError(
                        f"Column '{self.columns_to_read[0]}' not found in CSV file.")  # 如果指定的列不存在于 CSV 文件中，抛出 ValueError 异常

            content = '\n'.join(content)  # 将 content 列表中的字符串用换行符连接成一个大字符串

            # 提取源文件路径或者使用指定的来源列作为文档的来源
            source = (
                row.get(self.source_column, None)  # 获取指定的来源列的值，如果不存在，则使用 None
                if self.source_column is not None
                else self.file_path  # 如果未指定来源列，则使用文件路径作为来源
            )

            metadata = {"source": source, "row": i}  # 创建包含来源和行索引的元数据字典

            # 遍历其他元数据列，将其添加到 metadata 字典中
=======
                raise RuntimeError(f"Error loading {self.file_path}") from e
        except Exception as e:
            raise RuntimeError(f"Error loading {self.file_path}") from e

        return docs

    def __read_file(self, csvfile: TextIOWrapper) -> List[Document]:
        docs = []
        csv_reader = csv.DictReader(csvfile, **self.csv_args)  # type: ignore
        for i, row in enumerate(csv_reader):
            content = []
            for col in self.columns_to_read:
                if col in row:
                    content.append(f'{col}:{str(row[col])}')
                else:
                    raise ValueError(f"Column '{self.columns_to_read[0]}' not found in CSV file.")
            content = '\n'.join(content)
            # Extract the source if available
            source = (
                row.get(self.source_column, None)
                if self.source_column is not None
                else self.file_path
            )
            metadata = {"source": source, "row": i}

>>>>>>> d45db7c71cc1d7c6f454aab8dc32da6b0299ee3d
            for col in self.metadata_columns:
                if col in row:
                    metadata[col] = row[col]

<<<<<<< HEAD
            # 创建 Document 对象，将内容和元数据传递给 Document 类的构造函数
            doc = Document(page_content=content, metadata=metadata)

            docs.append(doc)  # 将创建的 Document 对象添加到 docs 列表中

        return docs  # 返回包含所有 Document 对象的列表
=======
            doc = Document(page_content=content, metadata=metadata)
            docs.append(doc)

        return docs
>>>>>>> d45db7c71cc1d7c6f454aab8dc32da6b0299ee3d


class RapidOCRDocLoader(UnstructuredFileLoader):
    def _get_elements(self) -> List:
<<<<<<< HEAD
        from unstructured.partition.text import partition_text  # 导入文本分割函数

        def doc2text(filepath):
            from docx.table import _Cell, Table  # 导入表格相关类
=======
        from unstructured.partition.text import partition_text

        def doc2text(filepath):
            from docx.table import _Cell, Table
>>>>>>> d45db7c71cc1d7c6f454aab8dc32da6b0299ee3d
            from docx.oxml.table import CT_Tbl
            from docx.oxml.text.paragraph import CT_P
            from docx.text.paragraph import Paragraph
            from docx import Document, ImagePart
<<<<<<< HEAD
            from PIL import Image  # 导入图像处理库

            ocr = get_ocr()  # 获取 OCR 实例
            doc = Document(filepath)  # 使用 docx 库打开文档
            resp = ""  # 初始化用于存储文档内容的字符串
=======
            from PIL import Image

            ocr = get_ocr()
            doc = Document(filepath)
            resp = ""
>>>>>>> d45db7c71cc1d7c6f454aab8dc32da6b0299ee3d

            def iter_block_items(parent):
                from docx.document import Document
                if isinstance(parent, Document):
                    parent_elm = parent.element.body
                elif isinstance(parent, _Cell):
                    parent_elm = parent._tc
                else:
<<<<<<< HEAD
                    raise ValueError("RapidOCRDocLoader 解析失败")  # 如果无法识别父元素类型，抛出异常
=======
                    raise ValueError("RapidOCRDocLoader parse fail")
>>>>>>> d45db7c71cc1d7c6f454aab8dc32da6b0299ee3d

                for child in parent_elm.iterchildren():
                    if isinstance(child, CT_P):
                        yield Paragraph(child, parent)
                    elif isinstance(child, CT_Tbl):
                        yield Table(child, parent)

<<<<<<< HEAD
            # 使用 tqdm 创建进度条，用于显示文档处理进度
            b_unit = tqdm(
                total=len(doc.paragraphs) + len(doc.tables),
                desc="RapidOCRDocLoader 处理块索引: 0"
            )
            for i, block in enumerate(iter_block_items(doc)):
                b_unit.set_description(
                    f"RapidOCRDocLoader 处理块索引: {i}"
                )
=======
            b_unit = tqdm(
                total=len(doc.paragraphs)+len(doc.tables),
                desc="RapidOCRDocLoader block index: 0"
            )
            for i, block in enumerate(iter_block_items(doc)):
                b_unit.set_description(
                    "RapidOCRDocLoader  block index: {}".format(i))
>>>>>>> d45db7c71cc1d7c6f454aab8dc32da6b0299ee3d
                b_unit.refresh()
                if isinstance(block, Paragraph):
                    resp += block.text.strip() + "\n"
                    images = block._element.xpath('.//pic:pic')  # 获取所有图片
                    for image in images:
                        for img_id in image.xpath('.//a:blip/@r:embed'):  # 获取图片id
                            part = doc.part.related_parts[img_id]  # 根据图片id获取对应的图片
                            if isinstance(part, ImagePart):
                                image = Image.open(BytesIO(part._blob))
                                result, _ = ocr(np.array(image))
                                if result:
                                    ocr_result = [line[1] for line in result]
                                    resp += "\n".join(ocr_result)
                elif isinstance(block, Table):
                    for row in block.rows:
                        for cell in row.cells:
                            for paragraph in cell.paragraphs:
                                resp += paragraph.text.strip() + "\n"
                b_unit.update(1)
            return resp

        text = doc2text(self.file_path)
<<<<<<< HEAD
        return partition_text(text=text, **self.unstructured_kwargs)  # 使用指定的文本分割器对文本进行分割
=======
        return partition_text(text=text, **self.unstructured_kwargs)
>>>>>>> d45db7c71cc1d7c6f454aab8dc32da6b0299ee3d


class RapidOCRImageLoader(UnstructuredFileLoader):
    def _get_elements(self) -> List:
<<<<<<< HEAD
        from unstructured.partition.text import partition_text  # 导入文本分割函数

        def img2text(filepath):
            resp = ""  # 初始化用于存储 OCR 结果的字符串
            ocr = get_ocr()  # 获取 OCR 实例
            result, _ = ocr(filepath)  # 对图像文件进行 OCR 识别
            if result:
                ocr_result = [line[1] for line in result]  # 提取 OCR 结果中的文本内容
                resp += "\n".join(ocr_result)  # 将识别结果拼接为一个字符串
            return resp

        text = img2text(self.file_path)  # 使用 img2text 函数处理图像文件，获取文本内容
        return partition_text(text=text, **self.unstructured_kwargs)  # 使用指定的文本分割器对文本进行分割
=======
        from unstructured.partition.text import partition_text

        def img2text(filepath):
            resp = ""
            ocr = get_ocr()
            result, _ = ocr(filepath)
            if result:
                ocr_result = [line[1] for line in result]
                resp += "\n".join(ocr_result)
            return resp

        text = img2text(self.file_path)
        return partition_text(text=text, **self.unstructured_kwargs)
>>>>>>> d45db7c71cc1d7c6f454aab8dc32da6b0299ee3d


class RapidOCRPDFLoader(UnstructuredFileLoader):
    def _get_elements(self) -> List:
<<<<<<< HEAD
        import cv2  # 导入 OpenCV 库用于图像处理
        from unstructured.partition.text import partition_text  # 导入文本分割函数

        def rotate_img(img, angle):
            h, w = img.shape[:2]  # 获取图像的高度和宽度
            rotate_center = (w / 2, h / 2)  # 设置旋转中心为图像中心点

            # 获取旋转矩阵
            # 参数1为旋转中心点;
            # 参数2为旋转角度, 正值为逆时针旋转, 负值为顺时针旋转;
            # 参数3为各向同性的比例因子, 1.0 表示原图大小
            M = cv2.getRotationMatrix2D(rotate_center, angle, 1.0)

            # 计算旋转后图像的新边界
            new_w = int(h * np.abs(M[0, 1]) + w * np.abs(M[0, 0]))
            new_h = int(h * np.abs(M[0, 0]) + w * np.abs(M[0, 1]))

=======
        import cv2
        from unstructured.partition.text import partition_text

        def rotate_img(img, angle):
            h, w = img.shape[:2]
            rotate_center = (w / 2, h / 2)
            # 获取旋转矩阵
            # 参数1为旋转中心点;
            # 参数2为旋转角度,正值-逆时针旋转;负值-顺时针旋转
            # 参数3为各向同性的比例因子,1.0原图，2.0变成原来的2倍，0.5变成原来的0.5倍
            M = cv2.getRotationMatrix2D(rotate_center, angle, 1.0)
            # 计算图像新边界
            new_w = int(h * np.abs(M[0, 1]) + w * np.abs(M[0, 0]))
            new_h = int(h * np.abs(M[0, 0]) + w * np.abs(M[0, 1]))
>>>>>>> d45db7c71cc1d7c6f454aab8dc32da6b0299ee3d
            # 调整旋转矩阵以考虑平移
            M[0, 2] += (new_w - w) / 2
            M[1, 2] += (new_h - h) / 2

<<<<<<< HEAD
            # 应用旋转矩阵进行图像旋转
            rotated_img = cv2.warpAffine(img, M, (new_w, new_h))
            return rotated_img  # 返回旋转后的图像

        # 在此处读取 PDF 文件并进行处理，具体实现可能涉及 PDF 解析和图像处理的代码
        # 返回处理后的文本元素列表，使用指定的文本分割器进行文本分割
        # partition_text 函数用于将文本按段落分割，并根据需要进一步处理
        text_elements = partition_text(text="", **self.unstructured_kwargs)
        return text_elements  # 返回处理后的文本元素列表

        def pdf2text(filepath):
            import fitz  # 导入 pyMuPDF 中的 fitz 模块，用于处理 PDF
            from tqdm import tqdm  # 导入 tqdm 用于显示进度条

            ocr = get_ocr()  # 获取 OCR 实例
            doc = fitz.open(filepath)  # 打开 PDF 文档
            resp = ""  # 初始化文本响应变量

            # 使用 tqdm 显示处理页面进度
            b_unit = tqdm(total=doc.page_count, desc="RapidOCRPDFLoader context page index: 0")
            for i, page in enumerate(doc):
                b_unit.set_description(f"RapidOCRPDFLoader context page index: {i}")  # 更新进度条描述
                b_unit.refresh()  # 刷新进度条

                text = page.get_text("")  # 提取页面文本内容
                resp += text + "\n"  # 将提取的文本内容添加到响应变量中

                img_list = page.get_image_info(xrefs=True)  # 获取页面中的图像信息列表
                for img in img_list:
                    if xref := img.get("xref"):  # 获取图像的引用编号
                        bbox = img["bbox"]  # 获取图像边界框信息

                        # 检查图像尺寸是否超过设定的阈值，根据阈值决定是否进行 OCR 处理
                        if ((bbox[2] - bbox[0]) / page.rect.width < PDF_OCR_THRESHOLD[0]
                                or (bbox[3] - bbox[1]) / page.rect.height < PDF_OCR_THRESHOLD[1]):
                            continue  # 图像尺寸小于阈值，跳过该图像

                        pix = fitz.Pixmap(doc, xref)  # 获取图像的像素映射对象
                        if int(page.rotation) != 0:  # 如果页面有旋转角度，则对图像进行旋转
                            img_array = np.frombuffer(pix.samples, dtype=np.uint8).reshape(pix.height, pix.width, -1)
                            tmp_img = Image.fromarray(img_array)
                            ori_img = cv2.cvtColor(np.array(tmp_img), cv2.COLOR_RGB2BGR)
                            rot_img = rotate_img(img=ori_img, angle=360 - page.rotation)  # 调用旋转图像函数
=======
            rotated_img = cv2.warpAffine(img, M, (new_w, new_h))
            return rotated_img

        def pdf2text(filepath):
            import fitz  # pyMuPDF里面的fitz包，不要与pip install fitz混淆

            ocr = get_ocr()
            doc = fitz.open(filepath)
            resp = ""

            b_unit = tqdm(total=doc.page_count, desc="RapidOCRPDFLoader context page index: 0")
            for i, page in enumerate(doc):
                b_unit.set_description("RapidOCRPDFLoader context page index: {}".format(i))
                b_unit.refresh()
                text = page.get_text("")
                resp += text + "\n"

                img_list = page.get_image_info(xrefs=True)
                for img in img_list:
                    if xref := img.get("xref"):
                        bbox = img["bbox"]
                        # 检查图片尺寸是否超过设定的阈值
                        if ((bbox[2] - bbox[0]) / page.rect.width < PDF_OCR_THRESHOLD[0]
                                or (bbox[3] - bbox[1]) / page.rect.height < PDF_OCR_THRESHOLD[1]):
                            continue
                        pix = fitz.Pixmap(doc, xref)
                        if int(page.rotation) != 0:  # 如果Page有旋转角度，则旋转图片
                            img_array = np.frombuffer(pix.samples, dtype=np.uint8).reshape(pix.height, pix.width, -1)
                            tmp_img = Image.fromarray(img_array)
                            ori_img = cv2.cvtColor(np.array(tmp_img), cv2.COLOR_RGB2BGR)
                            rot_img = rotate_img(img=ori_img, angle=360 - page.rotation)
>>>>>>> d45db7c71cc1d7c6f454aab8dc32da6b0299ee3d
                            img_array = cv2.cvtColor(rot_img, cv2.COLOR_RGB2BGR)
                        else:
                            img_array = np.frombuffer(pix.samples, dtype=np.uint8).reshape(pix.height, pix.width, -1)

<<<<<<< HEAD
                        result, _ = ocr(img_array)  # 对图像进行 OCR 识别
                        if result:
                            ocr_result = [line[1] for line in result]  # 提取 OCR 结果的文本内容
                            resp += "\n".join(ocr_result)  # 将 OCR 结果添加到响应变量中

                # 更新进度条
                b_unit.update(1)

            return resp  # 返回处理后的文本内容
=======
                        result, _ = ocr(img_array)
                        if result:
                            ocr_result = [line[1] for line in result]
                            resp += "\n".join(ocr_result)

                # 更新进度
                b_unit.update(1)
            return resp
>>>>>>> d45db7c71cc1d7c6f454aab8dc32da6b0299ee3d

        text = pdf2text(self.file_path)
        return partition_text(text=text, **self.unstructured_kwargs)


class OpenParserPDFLoader(UnstructuredFileLoader):
    def _get_elements(self) -> List:
<<<<<<< HEAD
        from unstructured.partition.text import partition_text  # 导入分割文本的函数

        def pdf2text(filepath):
            from openparse import DocumentParser  # 导入 openparse 中的 DocumentParser 类

            parser = DocumentParser()  # 创建 DocumentParser 实例
            parsed_content = parser.parse(filepath)  # 解析 PDF 文档内容
            resp = "\n".join(node.text for node in parsed_content.nodes)  # 将解析结果中的文本内容连接成字符串
            return resp  # 返回解析后的文本内容字符串

        text = pdf2text(self.file_path)  # 调用 pdf2text 函数处理当前实例的 PDF 文件路径
        return partition_text(text=text, **self.unstructured_kwargs)  # 使用指定的文本分割函数对文本进行分割，并返回结果列表
=======
        from unstructured.partition.text import partition_text

        def pdf2text(filepath):
            from openparse import DocumentParser

            parser = DocumentParser()
            parsed_content = parser.parse(filepath)
            resp = "\n".join(node.text for node in parsed_content.nodes)
            return resp

        text = pdf2text(self.file_path)
        return partition_text(text=text, **self.unstructured_kwargs)
>>>>>>> d45db7c71cc1d7c6f454aab8dc32da6b0299ee3d


class RapidOCRPPTLoader(UnstructuredFileLoader):
    def _get_elements(self) -> List:
<<<<<<< HEAD
        from unstructured.partition.text import partition_text  # 导入分割文本的函数

        def ppt2text(filepath):
            from pptx import Presentation  # 导入 pptx 库中的 Presentation 类

            ocr = get_ocr()  # 获取 OCR 实例
            prs = Presentation(filepath)  # 使用 pptx 打开指定路径的 PPT 文件
            resp = ""  # 初始化响应文本字符串

            def extract_text(shape):
                nonlocal resp  # 使用外部函数中的 resp 变量
                if shape.has_text_frame:
                    resp += shape.text.strip() + "\n"  # 如果形状包含文本框，则提取并添加文本内容到 resp
=======
        from unstructured.partition.text import partition_text

        def ppt2text(filepath):
            from pptx import Presentation

            ocr = get_ocr()
            prs = Presentation(filepath)
            resp = ""

            def extract_text(shape):
                nonlocal resp
                if shape.has_text_frame:
                    resp += shape.text.strip() + "\n"
>>>>>>> d45db7c71cc1d7c6f454aab8dc32da6b0299ee3d
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            for paragraph in cell.text_frame.paragraphs:
<<<<<<< HEAD
                                resp += paragraph.text.strip() + "\n"  # 如果形状是表格，则提取表格中的每个段落的文本内容
                if shape.shape_type == 13:  # 如果形状是图片
                    image = Image.open(BytesIO(shape.image.blob))  # 使用 BytesIO 打开图片数据
                    result, _ = ocr(np.array(image))  # 使用 OCR 进行图片识别
                    if result:
                        ocr_result = [line[1] for line in result]  # 提取 OCR 结果中的文本内容
                        resp += "\n".join(ocr_result)  # 将 OCR 结果添加到 resp 中
                elif shape.shape_type == 6:  # 如果形状是组合
                    for child_shape in shape.shapes:
                        extract_text(child_shape)  # 递归提取子形状中的文本内容

            b_unit = tqdm(total=len(prs.slides), desc="RapidOCRPPTLoader slide index: 1")  # 创建进度条
=======
                                resp += paragraph.text.strip() + "\n"
                if shape.shape_type == 13:  # 13 表示图片
                    image = Image.open(BytesIO(shape.image.blob))
                    result, _ = ocr(np.array(image))
                    if result:
                        ocr_result = [line[1] for line in result]
                        resp += "\n".join(ocr_result)
                elif shape.shape_type == 6:  # 6 表示组合
                    for child_shape in shape.shapes:
                        extract_text(child_shape)

            b_unit = tqdm(total=len(prs.slides), desc="RapidOCRPPTLoader slide index: 1")
>>>>>>> d45db7c71cc1d7c6f454aab8dc32da6b0299ee3d
            # 遍历所有幻灯片
            for slide_number, slide in enumerate(prs.slides, start=1):
                b_unit.set_description(
                    "RapidOCRPPTLoader slide index: {}".format(slide_number)
                )
                b_unit.refresh()
                sorted_shapes = sorted(
                    slide.shapes, key=lambda x: (x.top, x.left)
<<<<<<< HEAD
                )  # 按形状的位置从上到下、从左到右排序
                for shape in sorted_shapes:
                    extract_text(shape)  # 提取每个形状中的文本内容
                b_unit.update(1)  # 更新进度条

            return resp  # 返回提取的所有文本内容字符串

        text = ppt2text(self.file_path)  # 调用 ppt2text 函数处理当前实例的 PPT 文件路径
        return partition_text(text=text, **self.unstructured_kwargs)  # 使用指定的文本分割函数对文本进行分割处理，并返回结果列表

=======
                )  # 从上到下、从左到右遍历
                for shape in sorted_shapes:
                    extract_text(shape)
                b_unit.update(1)
            return resp

        text = ppt2text(self.file_path)
        return partition_text(text=text, **self.unstructured_kwargs)
>>>>>>> d45db7c71cc1d7c6f454aab8dc32da6b0299ee3d
